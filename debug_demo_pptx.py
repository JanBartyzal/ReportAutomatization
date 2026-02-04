"""
Debug script to analyze the demo PowerPoint file and understand the table structure.
"""
import sys
sys.path.append("backend")

from pptx import Presentation
from app.services.pptx_service import PowerpointManager

# Load the demo presentation
demo_path = "docs/demo/DemoPage.pptx"
prs = Presentation(demo_path)

print(f"Total slides: {len(prs.slides)}\n")

for slide_idx, slide in enumerate(prs.slides):
    print(f"=" * 80)
    print(f"SLIDE {slide_idx + 1}")
    print(f"=" * 80)
    
    # Get title
    title = "Untitled"
    if slide.shapes.title and slide.shapes.title.text:
        title = slide.shapes.title.text
    print(f"Title: {title}\n")
    
    # Analyze shapes
    print(f"Total shapes: {len(slide.shapes)}\n")
    
    text_shapes = []
    
    for shape_idx, shape in enumerate(slide.shapes):
        print(f"Shape {shape_idx}:")
        print(f"  Type: {shape.shape_type}")
        print(f"  Has table: {shape.has_table if hasattr(shape, 'has_table') else 'N/A'}")
        print(f"  Has text frame: {shape.has_text_frame if hasattr(shape, 'has_text_frame') else 'N/A'}")
        
        if hasattr(shape, 'has_text_frame') and shape.has_text_frame and shape.text:
            text_preview = shape.text[:100].replace('\n', ' ')
            print(f"  Text: {text_preview}")
            print(f"  Position: top={shape.top}, left={shape.left}, width={shape.width}, height={shape.height}")
            
            text_shapes.append({
                "index": shape_idx,
                "text": shape.text,
                "top": shape.top,
                "left": shape.left,
                "width": shape.width,
                "height": shape.height
            })
        
        print()
    
    # Sort text shapes by position
    print(f"\nText shapes sorted by position (top, left):")
    print(f"=" * 80)
    sorted_shapes = sorted(text_shapes, key=lambda s: (s["top"], s["left"]))
    
    for s in sorted_shapes:
        print(f"Idx {s['index']:2d} | Top: {s['top']:8.0f} | Left: {s['left']:8.0f} | Text: {s['text'][:60].replace(chr(10), ' ')}")
    
    print()

# Now test extraction with PowerpointManager
print(f"\n" + "=" * 80)
print(f"TESTING EXTRACTION WITH PowerpointManager")
print(f"=" * 80)

manager = PowerpointManager()
prs = manager.load_powerpoint(demo_path)
slides = manager.extract_slides(prs)

for slide_data in slides:
    print(f"\nSlide {slide_data.slide_index}: {slide_data.title}")
    print(f"  Table data rows: {len(slide_data.table_data)}")
    print(f"  Pseudo tables: {len(slide_data.pseudo_tables)}")
    print(f"  Images: {len(slide_data.image_data)}")
    print(f"  Text elements: {len(slide_data.text_content)}")
    
    if slide_data.pseudo_tables:
        for idx, pt in enumerate(slide_data.pseudo_tables):
            print(f"\n  Pseudo Table {idx}:")
            print(f"    Confidence: {pt.get('confidence_score')}")
            print(f"    Data ({len(pt.get('data', []))} rows x {len(pt.get('data', [[]])[0]) if pt.get('data') else 0} cols):")
            for row_idx, row in enumerate(pt.get('data', [])):
                print(f"      Row {row_idx}: {row}")
    
    if slide_data.table_data:
        print(f"\n  Native table data:")
        for row_idx, row in enumerate(slide_data.table_data[:5]):  # Show first 5 rows
            print(f"    Row {row_idx}: {row}")

print("\nDebug complete!")
