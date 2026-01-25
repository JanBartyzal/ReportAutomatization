"""
PowerPoint file processing and table extraction.

This module handles reading PPTX files, extracting tables from various sources
(native tables, Excel OLE objects, images), and creating PowerPoint presentations
from structured data.
"""

import logging
import io
import base64
from typing import List, Dict, Any, Optional
import pandas as pd
from pptx import Presentation
from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.slide import Slide
from pptx.shapes.base import BaseShape
from models import SlideData
from table_data import TableDataProcessor
from images_tesseract import TableImageData

logger = logging.getLogger(__name__)


class PowerpointManager:
    """
    Manager for PowerPoint file operations.
    
    Handles both creation and extraction of PowerPoint presentations,
    with intelligent table detection from multiple sources.
    """
    
    def __init__(self) -> None:
        """Initialize PowerPoint manager with default template and processors."""
        self.template_path = "template.pptx"
        self.image_processor = TableImageData()

    def create_powerpoint(self, slide_data: List[SlideData]) -> Presentation:
        """
        Create PowerPoint presentation from slide data.
        
        Args:
            slide_data: List of SlideData objects with content
            
        Returns:
            Presentation object ready to be saved
        """
        prs = Presentation(self.template_path)
        for slide in slide_data:
            # Layout with header and body (usually index 1)
            try:
                slide_layout = prs.slide_layouts[1]
            except IndexError:
                slide_layout = prs.slide_layouts[0]  # Fallback
                
            slide_obj = prs.slides.add_slide(slide_layout)
            if slide_obj.shapes.title:
                slide_obj.shapes.title.text = slide.title
                
            # Create table if data exists
            if slide.table_data:
                rows = len(slide.table_data)
                cols = 2  # Simplified default
                if rows > 0:
                     # Try to detect cols from first row keys
                     cols = len(slide.table_data[0].keys())

                table = slide_obj.shapes.add_table(
                    rows=rows+1,
                    cols=cols,
                    left=Inches(1),
                    top=Inches(2),
                    width=Inches(8),
                    height=Inches(5)
                ).table
                
                # Header
                headers = list(slide.table_data[0].keys()) if rows > 0 else ["Key", "Value"]
                for j, h in enumerate(headers):
                    table.cell(0, j).text = str(h)
                
                # Data
                for i, row in enumerate(slide.table_data):
                    for j, h in enumerate(headers):
                        table.cell(i+1, j).text = str(row.get(h, ""))
                        
        return prs

    def save_powerpoint(self, prs: Presentation, output_path: str) -> None:
        """
        Save PowerPoint presentation to file.
        
        Args:
            prs: Presentation object to save
            output_path: File path for output PPTX file
        """
        prs.save(output_path)

    def load_powerpoint(self, input_path: str) -> Presentation:
        """
        Load PowerPoint presentation from file.
        
        Args:
            input_path: Path to PPTX file
            
        Returns:
            Loaded Presentation object
        """
        return Presentation(input_path)

    def extract_native_table(self, shape: BaseShape) -> Optional[List[Dict[str, Any]]]:
        """
        Extract data from a native PowerPoint table shape.
        
        Converts table to list of dictionaries using first row as headers.
        
        Args:
            shape: PowerPoint shape object
            
        Returns:
            List of row dictionaries, or None if not a table or extraction fails
        """
        table_data = []
        if not shape.has_table:
            return None
            
        for row in shape.table.rows:
            row_data = []
            for cell in row.cells:
                row_data.append(cell.text)
            table_data.append(row_data)
            
        # Convert to list of dicts (assuming first row is header)
        if len(table_data) > 1:
            headers = table_data[0]
            result = []
            for row in table_data[1:]:
                # Zip headers with row data, handling potential length mismatches safely
                item = {}
                for i in range(len(headers)):
                    key = headers[i]
                    val = row[i] if i < len(row) else ""
                    item[key] = val
                result.append(item)
            return result
        return None

    def extract_excel_object(self, shape: BaseShape) -> Optional[List[Dict[str, Any]]]:
        """
        Extract data from an embedded Excel OLE object.
        
        Attempts to read Excel data from OLE object blob and convert to
        list of dictionaries.
        
        Args:
            shape: PowerPoint shape object
            
        Returns:
            List of row dictionaries, or None if not Excel or extraction fails
        """
        try:
            if shape.shape_type == MSO_SHAPE_TYPE.EMBEDDED_OLE_OBJECT:
                # Check prog_id for Excel
                if hasattr(shape,  'ole_format') and shape.ole_format.prog_id and \
                   ("excel" in shape.ole_format.prog_id.lower() or "worksheet" in shape.ole_format.prog_id.lower()):
                    
                    excel_blob = shape.ole_format.blob
                    try:
                        # python-pptx extraction of OLE object blob
                        # Sometimes this is the raw OLE container, sometimes just the file.
                        # pd.read_excel usually handles .xlsx or .xls bytes.
                        df = pd.read_excel(io.BytesIO(excel_blob))
                        # Basic cleanup: replace NaN with None or empty string for JSON compatibility
                        df = df.fillna("")
                        return df.to_dict(orient='records')
                    except Exception as e:
                        logger.error(f"Error reading Excel blob: {e}")
                        return None
        except Exception as e:
            logger.error(f"Error checking OLE object: {e}")
            return None
        return None

    def extract_text_from_slide(self, slide: Slide) -> List[str]:
        """
        Extract text content from all shapes in the slide that have a text frame.
        
        Args:
            slide: PowerPoint slide object
            
        Returns:
            List of text strings from slide shapes
        """
        text_content = []
        for shape in slide.shapes:
            if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                if shape.text and shape.text.strip():
                    text_content.append(shape.text.strip())
        return text_content

    def extract_slides(self, prs: Presentation) -> List[SlideData]:
        """
        Extract all slides with tables, images, and text from presentation.
        
        Implements multi-tier extraction strategy:
        1. Native PowerPoint tables
        2. Embedded Excel OLE objects
        3. Image-based tables (via OCR)
        
        Args:
            prs: Presentation object to extract from
            
        Returns:
            List of SlideData objects with extracted content
        """
        slides = []
        
        for i, slide in enumerate(prs.slides):
            extracted_images = []
            slide_tables_data = []  # Store all extracted data (native, excel)

            title = "Untitled"
            if slide.shapes.title and slide.shapes.title.text:
                title = slide.shapes.title.text
            
            # Extract text content
            text_content = self.extract_text_from_slide(slide)

            for shape in slide.shapes:
                # 1. Native Table
                if shape.has_table:
                    data = self.extract_native_table(shape)
                    if data:
                        slide_tables_data.extend(data)
                        continue  # Processed as table
                
                # 2. Excel OLE Object
                if shape.shape_type == MSO_SHAPE_TYPE.EMBEDDED_OLE_OBJECT:
                    data = self.extract_excel_object(shape)
                    if data:
                        slide_tables_data.extend(data)
                        continue  # Processed as Excel

                # 3. Image (Fallthrough)
                if hasattr(shape, "image"):
                    image_blob = shape.image.blob
                    local_table_data = self.image_processor.smart_extract(image_blob)
                    logger.debug(f"Image extraction result: {local_table_data}")
                    if local_table_data["data"]:
                        logger.info(f"Extracted table data from image: {local_table_data}")
                        slide_tables_data.extend(local_table_data["data"])
                    else:
                        # Convert to base64 for JSON transport
                        b64_img = base64.b64encode(image_blob).decode('utf-8')
                        extracted_images.append({
                        "slide_index": i + 1,
                        "image_base64": b64_img
                    })
                    
            slide_data = SlideData(
                slide_index=i + 1,
                title=title,
                table_data=slide_tables_data,
                image_data=extracted_images,
                text_content=text_content
            )
            slides.append(slide_data)
        return slides
    
    def get_image_from_slide(self, slide_data: SlideData, image_index: int) -> Optional[bytes]:
        """
        Get image bytes from slide data by index.
        
        Args:
            slide_data: SlideData object containing images
            image_index: Index of image to retrieve
            
        Returns:
            Decoded image bytes, or None if index out of range
        """
        if image_index < len(slide_data.image_data):
            return base64.b64decode(slide_data.image_data[image_index]['image_base64'])
        return None
    
    def extract_data(self, slide_data: SlideData) -> List[Dict[str, Any]]:
        """
        Unified extraction method for slide data.
        
        Priority:
        - If we already extracted native/excel tables, return them.
        - If not, try to extract data from the first image.
        
        Args:
            slide_data: SlideData object to extract from
            
        Returns:
            List of extracted table rows as dictionaries
        """
        if slide_data.table_data:
            logger.info(f"Extracting data from NATIVE/EXCEL source for slide {slide_data.slide_index}")
            return slide_data.table_data
            
        # Fallback to Image processing
        if slide_data.image_data:
            logger.info(f"Extracting data from IMAGE source for slide {slide_data.slide_index}")
            image_bytes = self.get_image_from_slide(slide_data, 0)
            if image_bytes is not None:
                # TableDataProcessor.extract_data returns JSON/Dict
                return TableDataProcessor().extract_data(image_bytes)
        
        return []