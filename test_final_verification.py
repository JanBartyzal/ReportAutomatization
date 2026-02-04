"""
FINAL VERIFICATION: Complete end-to-end test with real demo PPTX file.
This will test the full extraction pipeline including the enhanced PseudoTableParser.
"""
import sys
import os
import json
sys.path.insert(0, os.path.join(os.path.dirname(__file__), 'backend'))

from pptx import Presentation

print("=" * 80)
print("FINAL VERIFICATION TEST")
print("=" * 80)

# Extract shapes manually from demo file to verify column-based detection
demo_path = "docs/demo/DemoPage.pptx"
prs = Presentation(demo_path)

for slide_idx, slide in enumerate(prs.slides):
    print(f"\nSlide {slide_idx + 1}: {slide.shapes.title.text if slide.shapes.title else 'Untitled'}")
    
    # Collect text shapes
    shapes = []
    for shape in slide.shapes:
        if hasattr(shape, 'has_text_frame') and shape.has_text_frame and shape.text and shape.text.strip():
            shapes.append({
                "text": shape.text.strip(),
                "top": shape.top,
                "left": shape.left,
                "width": shape.width,
                "height": shape.height
            })
    
    print(f"  Found {len(shapes)} text shapes")
    
    # Test pseudo-table detection
    from app.services.parsers.ppt_shapes import PseudoTableParser
    
    parser = PseudoTableParser()  # Use default params (now includes proper EMU tolerances)
    results = parser.parse(shapes)
    
    print(f"  Detected {len(results)} pseudo-table(s)")
    
    if results:
        for idx, table in enumerate(results):
            data = table.get('data', [])
            print(f"\n  Pseudo-Table {idx}:")
            print(f"    Confidence: {table.get('confidence_score')}")
            print(f"    Dimensions: {len(data)} rows x {len(data[0]) if data else 0} cols")
            
            # Show table structure
            print(f"\n    Table Content:")
            for row_idx, row in enumerate(data[:8]):  # Show first 8 rows
                row_str = " | ".join([str(cell)[:30] if cell else "None" for cell in row])
                print(f"      Row {row_idx}: {row_str}")
            
            if len(data) > 8:
                print(f"      ... ({len(data) - 8} more rows)")
    
    # Write results to JSON
    output = {
        "slide_index": slide_idx + 1,
        "shapes_count": len(shapes),
        "pseudo_tables_detected": len(results),
        "results": results
    }
    
    with open('final_verification.json', 'w', encoding='utf-8') as f:
        json.dump(output, f, indent=2, ensure_ascii=False)
    
    print(f"\n  Full results saved to: final_verification.json")

print("\n" + "=" * 80)
print("VERIFICATION COMPLETE")
print("=" * 80)
