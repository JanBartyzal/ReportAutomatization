"""
Simple script to analyze the demo PowerPoint file structure.
"""
from pptx import Presentation
import json

# Load the demo presentation
demo_path = "docs/demo/DemoPage.pptx"
prs = Presentation(demo_path)

output = []
output.append(f"Total slides: {len(prs.slides)}\n")

for slide_idx, slide in enumerate(prs.slides):
    output.append(f"=" * 80)
    output.append(f"SLIDE {slide_idx + 1}")
    output.append(f"=" * 80)
    
    # Get title
    title = "Untitled"
    if slide.shapes.title and slide.shapes.title.text:
        title = slide.shapes.title.text
    output.append(f"Title: {title}\n")
    
    # Analyze shapes
    output.append(f"Total shapes: {len(slide.shapes)}\n")
    
    text_shapes = []
    
    for shape_idx, shape in enumerate(slide.shapes):
        has_table = shape.has_table if hasattr(shape, 'has_table') else False
        has_text_frame = shape.has_text_frame if hasattr(shape, 'has_text_frame') else False
        
        output.append(f"Shape {shape_idx}: Type={shape.shape_type}, Table={has_table}, TextFrame={has_text_frame}")
        
        if has_text_frame and shape.text:
            text_preview = shape.text[:100].replace('\n', ' ')
            output.append(f"  Text: {text_preview}")
            output.append(f"  Position: top={shape.top}, left={shape.left}, width={shape.width}, height={shape.height}")
            
            text_shapes.append({
                "index": shape_idx,
                "text": shape.text,
                "top": shape.top,
                "left": shape.left,
                "width": shape.width,
                "height": shape.height
            })
        
        output.append("")
    
    # Sort text shapes by position
    output.append(f"\nText shapes sorted by position (top, left):")
    output.append(f"=" * 80)
    sorted_shapes = sorted(text_shapes, key=lambda s: (s["top"], s["left"]))
    
    for s in sorted_shapes:
        text_clean = s['text'].replace('\n', ' | ')
        output.append(f"Idx {s['index']:2d} | Top: {s['top']:8.0f} | Left: {s['left']:8.0f} | Text: {text_clean[:80]}")
    
    output.append("")

output.append("\nAnalysis complete!")

# Write to file as UTF-8
with open('demo_analysis_utf8.txt', 'w', encoding='utf-8') as f:
    f.write('\n'.join(output))

print('\n'.join(output[:50]))  # Print first 50 lines
print(f"\n... Full output written to demo_analysis_utf8.txt ({len(output)} lines)")
