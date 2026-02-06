"""
Targeted test: Extract only the specific shapes that form the table.
Based on our earlier analysis, we know which shapes should form the table.
"""
import sys
import os
import json
sys.path.insert(0, os.path.join(os.path.dirname(__file__), 'backend'))

from pptx import Presentation
from app.services.parsers.ppt_shapes import PseudoTableParser

demo_path = "docs/demo/DemoPage.pptx"
prs = Presentation(demo_path)
slide = prs.slides[0]

# Define the indices of shapes that form the table based on our analysis
# Headers: 22, 23, 27, 28
# Column 1 (Category): 12-19  
# Column 2 (Descriptions): 60, 59, 58, 55, 53, 54, 52, 56
# Column 3 (Net savings): 67, 66, 65, 61, 62, 63, 64, 57
# Column 4 (Additional savings): 30-37

table_shape_indices = [
    # Headers
    22, 23, 27, 28,
    # Row data - first 5 rows only for cleaner test
    12, 60, 67, 30,  # Row 1
    13, 59, 66, 31,  # Row 2
    14, 58, 65, 32,  # Row 3
    15, 55, 61, 33,  # Row 4
    16, 53, 62, 34,  # Row 5
]

shapes = []
for idx in table_shape_indices:
    shape = slide.shapes[idx]
    if hasattr(shape, 'has_text_frame') and shape.has_text_frame and shape.text:
        shapes.append({
            "text": shape.text.strip(),
            "top": shape.top,
            "left": shape.left,
            "width": shape.width,
            "height": shape.height
        })

print(f"Testing with {len(shapes)} carefully selected table shapes")

parser = PseudoTableParser()  # Default tolerances
results = parser.parse(shapes)

print(f"Detected: {len(results)} table(s)\n")

if results:
    for idx, table in enumerate(results):
        data = table.get('data', [])
        print(f"Table {idx}:")
        print(f"  Confidence: {table.get('confidence_score')}")
        print(f"  Dimensions: {len(data)} rows x {len(data[0]) if data else 0} cols\n")
        
        print("  Content:")
        for row_idx, row in enumerate(data):
            row_str = " | ".join([str(cell)[:25] if cell else "None" for cell in row])
            print(f"    Row {row_idx}: {row_str}")
            
        # Save to JSON
        with open('targeted_test.json', 'w', encoding='utf-8') as f:
            json.dump({"success": True, "table": table}, f, indent=2, ensure_ascii=False)
        print("\n  Saved to targeted_test.json")
else:
    print("No tables detected")
    print("\nDEBUG INFO:")
    print(f"  Parser config: x_tol={parser.x_tolerance}, y_tol={parser.y_tolerance}")
    print(f"  Min requirements: {parser.min_rows} rows x {parser.min_cols} cols")
    print(f"  Input shapes: {len(shapes)}")
    
    with open('targeted_test.json', 'w', encoding='utf-8') as f:
        json.dump({
            "success": False,
            "shapes_count": len(shapes),
            "shapes": shapes[:5]  # First 5 for inspection
        }, f, indent=2, ensure_ascii=False)
