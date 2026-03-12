"""Tests for table_generator module."""

from __future__ import annotations

import io

import pytest
from pptx import Presentation
from pptx.util import Inches

from src.generators.pptx.service.table_generator import populate_table


def _make_slide_with_table(name: str, cols: int = 3, rows: int = 2) -> tuple:
    """Create a presentation with a table shape and return (prs, slide)."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    table_shape = slide.shapes.add_table(rows, cols, Inches(1), Inches(1), Inches(6), Inches(2))
    table_shape.name = name
    return prs, slide


class TestPopulateTable:
    """Test populate_table function."""

    def test_populate_basic_table(self):
        prs, slide = _make_slide_with_table("DataTable", cols=3, rows=2)

        result = populate_table(
            slide,
            "DataTable",
            headers=["Name", "Value", "Unit"],
            rows=[["CPU", "85", "%"]],
        )

        assert result is True

        # Verify header was set
        table = None
        for shape in slide.shapes:
            if shape.name == "DataTable":
                table = shape.table
                break
        assert table is not None
        assert table.rows[0].cells[0].text_frame.paragraphs[0].runs[0].text == "Name"

    def test_shape_not_found(self):
        prs, slide = _make_slide_with_table("DataTable")
        result = populate_table(slide, "NonExistent", ["A"], [["1"]])
        assert result is False

    def test_shape_not_a_table(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(1))
        txBox.name = "NotATable"
        txBox.text_frame.text = "text"

        result = populate_table(slide, "NotATable", ["A"], [["1"]])
        assert result is False

    def test_multiple_data_rows(self):
        prs, slide = _make_slide_with_table("DataTable", cols=2, rows=2)

        rows_data = [
            ["Server A", "99.9"],
            ["Server B", "95.2"],
            ["Server C", "88.0"],
        ]

        result = populate_table(
            slide,
            "DataTable",
            headers=["Server", "Uptime %"],
            rows=rows_data,
        )

        assert result is True
