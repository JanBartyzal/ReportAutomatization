"""Tests for pptx_renderer module."""

from __future__ import annotations

import io

from pptx import Presentation
from pptx.util import Inches

from src.generators.pptx.service.pptx_renderer import RenderResult, render, ChartInput
from src.generators.pptx.service.chart_generator import ChartSeriesData


def _make_template_with_text_placeholder(tag: str = "{{company_name}}") -> bytes:
    """Create a minimal PPTX template with a text placeholder."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    txBox.name = "TextBox 1"
    tf = txBox.text_frame
    tf.text = f"Company: {tag}"

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.getvalue()


def _make_template_with_multiple_placeholders() -> bytes:
    """Create a PPTX template with text, table, and chart placeholders."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Text placeholder
    txBox = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(4), Inches(0.5))
    txBox.name = "Title"
    txBox.text_frame.text = "Report for {{company_name}}"

    # Chart placeholder in a separate text box
    chartBox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(6), Inches(3))
    chartBox.name = "ChartArea"
    chartBox.text_frame.text = "{{CHART:revenue}}"

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.getvalue()


class TestRender:
    """Test the render function."""

    def test_text_substitution(self):
        template = _make_template_with_text_placeholder()
        result = render(
            template_bytes=template,
            text_data={"company_name": "Acme Corp"},
            table_data={},
            chart_data={},
        )

        assert isinstance(result, RenderResult)
        assert len(result.pptx_bytes) > 0
        assert result.missing_placeholders == []

        # Verify the text was replaced
        prs = Presentation(io.BytesIO(result.pptx_bytes))
        slide = prs.slides[0]
        texts = [shape.text_frame.text for shape in slide.shapes if shape.has_text_frame]
        assert any("Acme Corp" in t for t in texts)
        assert not any("{{company_name}}" in t for t in texts)

    def test_missing_text_placeholder(self):
        template = _make_template_with_text_placeholder()
        result = render(
            template_bytes=template,
            text_data={},  # No data provided
            table_data={},
            chart_data={},
        )

        assert len(result.missing_placeholders) == 1
        assert "{{company_name}}" in result.missing_placeholders

        # Verify DATA MISSING text
        prs = Presentation(io.BytesIO(result.pptx_bytes))
        slide = prs.slides[0]
        texts = [shape.text_frame.text for shape in slide.shapes if shape.has_text_frame]
        assert any("DATA MISSING" in t for t in texts)

    def test_chart_insertion(self):
        template = _make_template_with_multiple_placeholders()
        result = render(
            template_bytes=template,
            text_data={"company_name": "Test Co"},
            table_data={},
            chart_data={
                "revenue": ChartInput(
                    chart_type="BAR",
                    labels=["Q1", "Q2", "Q3"],
                    series=[ChartSeriesData(name="Revenue", values=[100.0, 200.0, 150.0])],
                ),
            },
        )

        assert len(result.pptx_bytes) > 0
        assert "{{CHART:revenue}}" not in result.missing_placeholders

    def test_output_is_valid_pptx(self):
        template = _make_template_with_text_placeholder()
        result = render(
            template_bytes=template,
            text_data={"company_name": "Valid Corp"},
            table_data={},
            chart_data={},
        )

        # Should be parseable as a valid PPTX
        prs = Presentation(io.BytesIO(result.pptx_bytes))
        assert len(prs.slides) == 1
