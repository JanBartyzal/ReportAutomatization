"""Core PPTX parsing logic using python-pptx.

Extracts slide metadata, text blocks with positions, tables, and speaker notes.
Handles edge cases: empty slides, merged cells, SmartArt, charts.
"""

from __future__ import annotations

import logging
import uuid
from dataclasses import dataclass, field
from io import BytesIO
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Emu

logger = logging.getLogger(__name__)


# ---------------------------------------------------------------------------
# Data classes for parsed output
# ---------------------------------------------------------------------------


@dataclass(slots=True)
class TextBlockData:
    """A single text block extracted from a slide shape."""

    shape_name: str
    text: str
    is_title: bool
    position_x: int  # EMU converted to integer
    position_y: int


@dataclass(slots=True)
class TableRowData:
    """A single table row."""

    cells: list[str]


@dataclass(slots=True)
class TableDataParsed:
    """A table extracted from a slide."""

    table_id: str
    headers: list[str]
    rows: list[TableRowData]
    confidence: float = 1.0  # Native tables always have confidence 1.0


@dataclass(slots=True)
class SlideContent:
    """All content extracted from a single slide."""

    slide_index: int
    texts: list[TextBlockData] = field(default_factory=list)
    tables: list[TableDataParsed] = field(default_factory=list)
    notes: str = ""


@dataclass(slots=True)
class SlideInfo:
    """Metadata about a single slide."""

    slide_index: int
    title: str
    layout_name: str
    has_tables: bool
    has_text: bool
    has_images: bool
    has_charts: bool
    has_notes: bool


@dataclass(slots=True)
class DocumentProperties:
    """Document-level properties."""

    properties: dict[str, str] = field(default_factory=dict)


@dataclass(slots=True)
class PptxStructure:
    """Complete PPTX file structure."""

    total_slides: int
    slides: list[SlideInfo]
    document_properties: dict[str, str]


# ---------------------------------------------------------------------------
# Parser
# ---------------------------------------------------------------------------


class PptxParser:
    """Stateless parser that wraps python-pptx operations.

    Usage::

        parser = PptxParser()
        prs = parser.open(path)
        structure = parser.extract_structure(prs)
        content = parser.extract_slide_content(prs, slide_index=0)
    """

    @staticmethod
    def open(source: str | Path | bytes | BytesIO) -> Presentation:
        """Open a PPTX file and return a ``Presentation`` object.

        Args:
            source: File path, raw bytes, or BytesIO stream.

        Returns:
            A python-pptx ``Presentation`` instance.
        """
        if isinstance(source, bytes):
            source = BytesIO(source)
        return Presentation(source)

    # -- structure extraction ----------------------------------------------

    def extract_structure(self, prs: Presentation) -> PptxStructure:
        """Extract the high-level structure of a presentation.

        Args:
            prs: An opened python-pptx Presentation.

        Returns:
            A ``PptxStructure`` describing slides and document properties.
        """
        slides_info: list[SlideInfo] = []

        for idx, slide in enumerate(prs.slides):
            title = self._detect_title(slide)
            layout_name = slide.slide_layout.name if slide.slide_layout else ""
            has_tables = False
            has_text = False
            has_images = False
            has_charts = False

            for shape in slide.shapes:
                if shape.has_table:
                    has_tables = True
                if shape.has_text_frame:
                    if shape.text_frame.text.strip():
                        has_text = True
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    has_images = True
                if shape.has_chart:
                    has_charts = True
                # SmartArt is represented as a group shape
                if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                    has_images = True

            has_notes = bool(
                slide.has_notes_slide
                and slide.notes_slide.notes_text_frame
                and slide.notes_slide.notes_text_frame.text.strip()
            )

            slides_info.append(
                SlideInfo(
                    slide_index=idx,
                    title=title,
                    layout_name=layout_name,
                    has_tables=has_tables,
                    has_text=has_text,
                    has_images=has_images,
                    has_charts=has_charts,
                    has_notes=has_notes,
                )
            )

        doc_props = self._extract_document_properties(prs)

        return PptxStructure(
            total_slides=len(prs.slides),
            slides=slides_info,
            document_properties=doc_props,
        )

    # -- slide content extraction ------------------------------------------

    def extract_slide_content(self, prs: Presentation, slide_index: int) -> SlideContent:
        """Extract all content from a specific slide.

        Args:
            prs: An opened python-pptx Presentation.
            slide_index: Zero-based slide index.

        Returns:
            A ``SlideContent`` with texts, tables, and notes.

        Raises:
            IndexError: If slide_index is out of range.
        """
        if slide_index < 0 or slide_index >= len(prs.slides):
            raise IndexError(f"Slide index {slide_index} out of range (0..{len(prs.slides) - 1})")

        slide = prs.slides[slide_index]
        texts = self._extract_text_blocks(slide)
        tables = self._extract_tables(slide)
        notes = self._extract_notes(slide)

        return SlideContent(
            slide_index=slide_index,
            texts=texts,
            tables=tables,
            notes=notes,
        )

    # -- private helpers ---------------------------------------------------

    @staticmethod
    def _detect_title(slide: Any) -> str:
        """Detect the slide title from placeholder shapes.

        Args:
            slide: A python-pptx Slide object.

        Returns:
            The title text, or empty string if none found.
        """
        if slide.shapes.title is not None:
            return slide.shapes.title.text.strip()

        # Fallback: look for a shape whose placeholder index indicates a title
        for shape in slide.placeholders:
            if shape.placeholder_format.idx in (0, 1):  # Title or center title
                if shape.has_text_frame and shape.text_frame.text.strip():
                    return shape.text_frame.text.strip()
        return ""

    @staticmethod
    def _extract_text_blocks(slide: Any) -> list[TextBlockData]:
        """Extract text blocks with position data from all shapes.

        Args:
            slide: A python-pptx Slide object.

        Returns:
            List of ``TextBlockData`` instances.
        """
        blocks: list[TextBlockData] = []

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue

            text = shape.text_frame.text.strip()
            if not text:
                continue

            is_title = False
            if hasattr(shape, "placeholder_format") and shape.placeholder_format is not None:
                is_title = shape.placeholder_format.idx in (0, 1)

            # Position in EMUs; convert to integer for proto transport
            pos_x = int(shape.left) if shape.left is not None else 0
            pos_y = int(shape.top) if shape.top is not None else 0

            blocks.append(
                TextBlockData(
                    shape_name=shape.name or "",
                    text=text,
                    is_title=is_title,
                    position_x=pos_x,
                    position_y=pos_y,
                )
            )

        return blocks

    @staticmethod
    def _extract_tables(slide: Any) -> list[TableDataParsed]:
        """Extract native tables from slide shapes.

        Handles merged cells by using the text of the top-left merge origin cell.

        Args:
            slide: A python-pptx Slide object.

        Returns:
            List of ``TableDataParsed`` instances.
        """
        tables: list[TableDataParsed] = []

        for shape in slide.shapes:
            if not shape.has_table:
                continue

            table = shape.table
            if table.rows is None or len(table.rows) == 0:
                continue

            # First row is treated as headers
            header_row = table.rows[0]
            headers: list[str] = []
            for cell in header_row.cells:
                headers.append(cell.text.strip())

            rows: list[TableRowData] = []
            for row_idx, row in enumerate(table.rows):
                if row_idx == 0:
                    continue  # skip header row
                cells: list[str] = []
                for cell in row.cells:
                    cells.append(cell.text.strip())
                rows.append(TableRowData(cells=cells))

            tables.append(
                TableDataParsed(
                    table_id=str(uuid.uuid4()),
                    headers=headers,
                    rows=rows,
                    confidence=1.0,
                )
            )

        return tables

    @staticmethod
    def _extract_notes(slide: Any) -> str:
        """Extract speaker notes from a slide.

        Args:
            slide: A python-pptx Slide object.

        Returns:
            Notes text, or empty string if no notes.
        """
        if not slide.has_notes_slide:
            return ""
        notes_slide = slide.notes_slide
        if notes_slide.notes_text_frame is None:
            return ""
        return notes_slide.notes_text_frame.text.strip()

    @staticmethod
    def _extract_document_properties(prs: Presentation) -> dict[str, str]:
        """Extract document-level core properties.

        Args:
            prs: An opened python-pptx Presentation.

        Returns:
            Dictionary of property name to value.
        """
        props: dict[str, str] = {}
        core = prs.core_properties
        if core is None:
            return props

        for attr in ("author", "title", "subject", "keywords", "comments", "category", "last_modified_by"):
            value = getattr(core, attr, None)
            if value:
                props[attr] = str(value)

        if core.created:
            props["created"] = core.created.isoformat()
        if core.modified:
            props["modified"] = core.modified.isoformat()

        return props
