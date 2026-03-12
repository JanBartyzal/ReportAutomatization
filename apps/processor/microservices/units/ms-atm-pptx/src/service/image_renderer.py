"""Slide image rendering using LibreOffice headless with python-pptx fallback.

Renders individual slides as 1280x720 PNG images.
"""

from __future__ import annotations

import asyncio
import logging
import shutil
import tempfile
from pathlib import Path
from typing import Final

from src.config import LIBREOFFICE_BIN, RENDER_HEIGHT, RENDER_WIDTH

logger = logging.getLogger(__name__)

_LIBREOFFICE_TIMEOUT: Final[int] = 60  # seconds


class ImageRenderer:
    """Render PPTX slides to PNG images.

    Primary strategy: LibreOffice headless conversion.
    Fallback: basic python-pptx shape rendering with Pillow (limited fidelity).

    Usage::

        renderer = ImageRenderer()
        png_bytes = await renderer.render_slide(pptx_path, slide_index=0)
    """

    def __init__(self) -> None:
        self._libreoffice_available: bool | None = None

    async def render_slide(self, pptx_path: str | Path, slide_index: int) -> bytes:
        """Render a single slide as a PNG image.

        Args:
            pptx_path: Path to the local PPTX file.
            slide_index: Zero-based index of the slide to render.

        Returns:
            PNG image bytes (1280x720).

        Raises:
            RuntimeError: If rendering fails with both strategies.
        """
        pptx_path = Path(pptx_path)

        if await self._is_libreoffice_available():
            try:
                return await self._render_with_libreoffice(pptx_path, slide_index)
            except Exception:
                logger.warning(
                    "LibreOffice rendering failed for slide %d, falling back to Pillow",
                    slide_index,
                    exc_info=True,
                )

        return self._render_fallback(pptx_path, slide_index)

    # -- LibreOffice rendering ---------------------------------------------

    async def _render_with_libreoffice(self, pptx_path: Path, slide_index: int) -> bytes:
        """Render via LibreOffice headless conversion.

        LibreOffice converts the entire PPTX to a directory of PNGs. We then
        select the specific slide image by index.

        Args:
            pptx_path: Path to the PPTX file.
            slide_index: Zero-based slide index.

        Returns:
            PNG bytes for the requested slide.
        """
        with tempfile.TemporaryDirectory(prefix="pptx_render_") as tmpdir:
            cmd = [
                LIBREOFFICE_BIN,
                "--headless",
                "--norestore",
                "--convert-to", "png",
                "--outdir", tmpdir,
                str(pptx_path),
            ]

            logger.debug("Running LibreOffice: %s", " ".join(cmd))

            proc = await asyncio.create_subprocess_exec(
                *cmd,
                stdout=asyncio.subprocess.PIPE,
                stderr=asyncio.subprocess.PIPE,
            )
            stdout, stderr = await asyncio.wait_for(proc.communicate(), timeout=_LIBREOFFICE_TIMEOUT)

            if proc.returncode != 0:
                raise RuntimeError(
                    f"LibreOffice exited with code {proc.returncode}: {stderr.decode(errors='replace')}"
                )

            # LibreOffice produces one PNG per slide (or a single PNG for the whole file).
            # File naming: <stem>.png or <stem>-<n>.png
            tmpdir_path = Path(tmpdir)
            png_files = sorted(tmpdir_path.glob("*.png"))

            if not png_files:
                raise RuntimeError("LibreOffice produced no PNG output")

            # If single file, all slides collapsed (unlikely but possible)
            if len(png_files) == 1:
                target = png_files[0]
            elif slide_index < len(png_files):
                target = png_files[slide_index]
            else:
                raise RuntimeError(
                    f"Slide index {slide_index} out of range; LibreOffice produced {len(png_files)} images"
                )

            raw_bytes = target.read_bytes()

            # Resize to target dimensions
            return self._resize_png(raw_bytes, RENDER_WIDTH, RENDER_HEIGHT)

    # -- Fallback rendering ------------------------------------------------

    @staticmethod
    def _render_fallback(pptx_path: Path, slide_index: int) -> bytes:
        """Basic fallback rendering using Pillow.

        Creates a white canvas with the slide title overlaid. This is a
        low-fidelity placeholder for when LibreOffice is not available.

        Args:
            pptx_path: Path to the PPTX file.
            slide_index: Zero-based slide index.

        Returns:
            PNG bytes for a simple rendered image.
        """
        from io import BytesIO

        from PIL import Image, ImageDraw, ImageFont
        from pptx import Presentation

        prs = Presentation(str(pptx_path))
        if slide_index < 0 or slide_index >= len(prs.slides):
            raise IndexError(f"Slide index {slide_index} out of range")

        slide = prs.slides[slide_index]
        title = ""
        if slide.shapes.title is not None:
            title = slide.shapes.title.text

        img = Image.new("RGB", (RENDER_WIDTH, RENDER_HEIGHT), color=(255, 255, 255))
        draw = ImageDraw.Draw(img)

        # Draw slide number and title
        label = f"Slide {slide_index + 1}"
        if title:
            label += f": {title}"

        try:
            font = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", 28)
        except (OSError, IOError):
            font = ImageFont.load_default()

        draw.text((40, 40), label, fill=(0, 0, 0), font=font)

        # Draw shape outlines for basic visual representation
        for shape in slide.shapes:
            if shape.left is None or shape.top is None:
                continue
            if shape.width is None or shape.height is None:
                continue

            # Scale from EMU to pixel coordinates
            slide_width = prs.slide_width or 9144000  # default 10 inches
            slide_height = prs.slide_height or 6858000  # default 7.5 inches

            x = int(shape.left / slide_width * RENDER_WIDTH)
            y = int(shape.top / slide_height * RENDER_HEIGHT)
            w = int(shape.width / slide_width * RENDER_WIDTH)
            h = int(shape.height / slide_height * RENDER_HEIGHT)

            draw.rectangle([x, y, x + w, y + h], outline=(200, 200, 200), width=1)

        buf = BytesIO()
        img.save(buf, format="PNG")
        return buf.getvalue()

    # -- Utilities ---------------------------------------------------------

    @staticmethod
    def _resize_png(data: bytes, width: int, height: int) -> bytes:
        """Resize PNG bytes to the target dimensions.

        Args:
            data: Original PNG bytes.
            width: Target width.
            height: Target height.

        Returns:
            Resized PNG bytes.
        """
        from io import BytesIO

        from PIL import Image

        img = Image.open(BytesIO(data))
        if img.size != (width, height):
            img = img.resize((width, height), Image.LANCZOS)

        buf = BytesIO()
        img.save(buf, format="PNG")
        return buf.getvalue()

    async def _is_libreoffice_available(self) -> bool:
        """Check whether LibreOffice is available on the system.

        The result is cached after the first check.

        Returns:
            True if LibreOffice can be invoked.
        """
        if self._libreoffice_available is not None:
            return self._libreoffice_available

        if shutil.which(LIBREOFFICE_BIN) is not None:
            self._libreoffice_available = True
        else:
            # Try running it (some containers have it on PATH but not via which)
            try:
                proc = await asyncio.create_subprocess_exec(
                    LIBREOFFICE_BIN, "--version",
                    stdout=asyncio.subprocess.PIPE,
                    stderr=asyncio.subprocess.PIPE,
                )
                await asyncio.wait_for(proc.communicate(), timeout=10)
                self._libreoffice_available = proc.returncode == 0
            except Exception:
                self._libreoffice_available = False

        if not self._libreoffice_available:
            logger.warning("LibreOffice not available; slide rendering will use fallback")

        return self._libreoffice_available
